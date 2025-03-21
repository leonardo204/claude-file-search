"""
PowerPoint (.pptx) 파일 핸들러

이 모듈은 PowerPoint 파일에서 텍스트를 추출하는 기능을 제공합니다.
"""

import logging
import traceback
from pathlib import Path
from typing import List, Dict, Any

from file_handlers.base import FileHandler

logger = logging.getLogger("file_search.pptx_handler")

class PPTXHandler(FileHandler):
    """PowerPoint 파일 핸들러"""
    
    def get_supported_extensions(self) -> List[str]:
        """
        지원하는 파일 확장자 목록을 반환합니다.
        
        Returns:
            지원하는 파일 확장자 목록
        """
        return ['.pptx']
    
    def get_type_description(self) -> str:
        """
        파일 형식에 대한 설명을 반환합니다.
        
        Returns:
            파일 형식 설명
        """
        return 'PowerPoint'
    
    def extract_text(self, file_path: Path, max_content_sections: int = 10) -> List[Dict[str, Any]]:
        """
        PowerPoint 파일에서 텍스트를 추출합니다.
        
        Args:
            file_path: 파일 경로
            max_content_sections: 처리할 최대 슬라이드 수
            
        Returns:
            슬라이드별 텍스트 정보 리스트
        """
        try:
            from pptx import Presentation
            
            logger.info(f"Extracting text from PowerPoint file: {file_path}")
            
            presentation = Presentation(file_path)
            slides = []
            
            # 최대 슬라이드 수 제한
            for i, slide in enumerate(presentation.slides):
                if i >= max_content_sections:
                    slides.append({
                        "section_number": i + 1,
                        "section_type": "slide",
                        "text": "... 추가 슬라이드 있음 (처리되지 않음) ..."
                    })
                    break
                    
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                slides.append({
                    "section_number": i + 1,
                    "section_type": "slide",
                    "text": slide_text.strip()
                })
            
            return slides
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint file {file_path}: {e}")
            logger.error(traceback.format_exc())
            return [{"section_number": 0, "section_type": "error", "text": f"PowerPoint 파일 처리 중 오류 발생: {str(e)}"}]